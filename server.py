import pptx
from pptx import Presentation
from pptx.util import Inches
from mcp.server import Server 
from mcp.types import Tool, TextContent 
import asyncio 
from typing import List, Tuple
from exa_py import Exa 
import os 

key = "38e47a00-6c20-4415-b7b3-5a08b2145df6" 
exa = Exa(api_key=key)

server = Server('powerpoint-presentation')

@server.list_tools()
async def list_tools()->List[Tool]:
    '''Tell the model what tools are available'''

    return [
        Tool(
            name='create_presentation',
            description='Creates a new PowerPoint Presentation with a title slide',
            inputSchema={
                'type':'object',
                'properties': {
                    'title':{
                        'type':'string',
                        'description': 'The title for the Presentation'
                    },
                    'slides': {
                        'type': "array",
                        'description': "Array of slidse to create",
                        'items': {
                            'type': 'object',
                            'properties':{
                                'title': {
                                    'type': 'string',
                                    'description': 'Title of the slide'
                                },
                                'content': {
                                    'type': 'string',
                                    'description': 'Main content/body text for the slide'
                                },
                                'bullet_points':{
                                    'type':'array',
                                    'description':'Optional buller points',
                                    'items': {'type': 'string'}
                                }
                            },
                            'required': ['title']
                        }
                    },
                    'filename':{
                        'type':'string',
                        'description':'Filename to save as (default: presentation.pptx',
                    },
                },
                "required":['title', 'slides']
            }
        ),
        Tool(
            name='search_web_presentation',
            description='Search web for information about a topic',
            inputSchema={
                'type':'object',
                'properties': {
                    'query':{
                        'type':'string',
                        'description':'The search query/topic to research'
                    },
                    'max_results':{
                        'type':'integer',
                        'description':'Maximum number of search results to return (default: 5)',
                        'default': 5
                    }
                },
                "required":['query']
            }
        )
    ]


@server.call_tool()
async def call_tool(name: str, arguments: dict)->List[TextContent]:
    '''Handle tool calls from the model'''
    if name == 'create_presentation':
        title_text = arguments['title']
        filename=arguments.get('filename', 'presentation.pptx')
        slides_data = arguments['slides']

        prs = Presentation()
        
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = title_text  
        
        for slide_data in slides_data:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_data['title']

            content_box = slide.placeholders[1]
            text_frame = content_box.text_frame 
            text_frame.clear()

            if 'content' in slide_data and slide_data['content']:
                p = text_frame.paragraphs[0]
                p.text = slide_data['content']
                p.level=0

            if 'bullet_points' in slide_data:
                for i, bullet in enumerate(slide_data['bullet_points']):
                    if i==0 and 'content' not in slide_data:
                        p=text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = bullet 
                    p.level= 1

        prs.save(filename)
        
        return [TextContent(
            type='text',
            text=f"Created presentation '{filename}' with title '{title_text}' and {len(slides_data)} slides"
        )]

    elif name == 'search_web_presentation':
        query = arguments['query']
        max_results = arguments.get('max_results', 5)
        try:
            results = exa.search_and_contents(
                query,
                num_results=max_results,
                text=True,
                type='auto',
            )
            summary = []
            for i, result in enumerate(results.results, 1):
                content = ''
                if getattr(result, 'text', None):
                    content = result.text[:500] + '...'

                summary.append(
                    f"{i}. {result.title}\n"
                    f"   {content}\n"
                    f"   Source: {result.url}\n"
                    f"   Published: {getattr(result, 'published_date', 'N/A')}\n"
                )
            search_sum = '\n'.join(summary)

            return [TextContent(
                type='text',
                text=f'Search results for "{query}":\n\n{search_sum}'
            )]

        except Exception as e:
            return [TextContent(
            type='text',
            text=f'Search failed: {str(e)}'
        )]

    raise ValueError(f"Unknown tool: {name}")

async def main():
    '''Run the server '''
    from mcp.server.stdio import stdio_server 
    async with stdio_server() as (read_stream, write_stream):
        await server.run(
            read_stream,
            write_stream,
            server.create_initialization_options()
        )

if __name__ == '__main__':
    asyncio.run(main())

